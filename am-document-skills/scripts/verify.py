#!/usr/bin/env python3
"""
Verify document output quality.
Usage: python3 verify.py <file_path>

Checks:
- File exists and not empty
- Format-specific validation
- Font presence (DOCX/PPTX)
- Vietnamese text rendering (PDF)
- Formula presence (XLSX)
"""

import sys
import os


# Vietnamese diacritics detection
VN_DIACRITICS = set('\u00e0\u00e1\u1ea3\u00e3\u1ea1\u0103\u1eaf\u1eb1\u1eb3\u1eb5\u1eb7\u00e2\u1ea5\u1ea7\u1ea9\u1eab\u1ead'
                    '\u00e8\u00e9\u1ebb\u1ebd\u1eb9\u00ea\u1ebf\u1ec1\u1ec3\u1ec5\u1ec7'
                    '\u00ec\u00ed\u1ec9\u0129\u1ecb'
                    '\u00f2\u00f3\u1ecf\u00f5\u1ecd\u00f4\u1ed1\u1ed3\u1ed5\u1ed7\u1ed9\u01a1\u1edb\u1edd\u1edf\u1ee1\u1ee3'
                    '\u00f9\u00fa\u1ee7\u0169\u1ee5\u01b0\u1ee9\u1eeb\u1eed\u1eef\u1ef1'
                    '\u1ef3\u00fd\u1ef7\u1ef9\u1ef5\u0111')

VN_UNDIACRITICED_WORDS = [
    'Tong', 'Quan', 'Phan', 'Tich', 'Bao', 'Cao', 'Doanh', 'Thu',
    'Nhan', 'Xet', 'Trien', 'Vong', 'Thanh', 'Thi', 'Nong', 'Thon',
    'Cong', 'Nghiep', 'Dich', 'Vu', 'Noi', 'Dung',
]


def check_vietnamese_diacritics(text):
    """Check if Vietnamese text has proper diacritics.
    Returns (has_vn_content, has_diacritics, warning)
    """
    if not text:
        return False, True, ''
    lower = text.lower()
    has_diacritics = any(c in lower for c in VN_DIACRITICS)
    # Check for common Vietnamese words without diacritics
    undiac_count = sum(1 for w in VN_UNDIACRITICED_WORDS if f' {w} ' in text or text.startswith(f'{w} '))
    if undiac_count >= 3 and not has_diacritics:
        return True, False, f'Vietnamese text appears to lack diacritics ({undiac_count} undiacriticed words found)'
    return has_diacritics, True, ''


def verify_docx(path):
    """Verify DOCX file."""
    from docx import Document
    results = []
    try:
        doc = Document(path)
        results.append(('File opens', True, ''))

        # Check page count (sections)
        results.append(('Has sections', len(doc.sections) > 0, f'{len(doc.sections)} sections'))

        # Check paragraphs
        para_count = len(doc.paragraphs)
        results.append(('Has content', para_count > 0, f'{para_count} paragraphs'))

        # Check headings
        headings = [p for p in doc.paragraphs if p.style.name.startswith('Heading')]
        results.append(('Has headings', len(headings) > 0, f'{len(headings)} headings'))

        # Check font
        fonts_used = set()
        for p in doc.paragraphs:
            for run in p.runs:
                if run.font.name:
                    fonts_used.add(run.font.name)
        has_be_vietnam = any('Vietnam' in f for f in fonts_used)
        has_arial = 'Arial' in fonts_used
        results.append(('Font set', True, ', '.join(fonts_used) or 'default'))
        if fonts_used and not has_be_vietnam and not has_arial:
            results.append(('Font warning', False, 'No Be Vietnam Pro or Arial found'))

        # Check tables
        tables = doc.tables
        results.append(('Tables', True, f'{len(tables)} tables'))

        # Check page margins (first section)
        s = doc.sections[0]
        margin_ok = s.top_margin is not None
        results.append(('Margins set', margin_ok, ''))

        # Check Vietnamese diacritics
        all_text = ' '.join(p.text for p in doc.paragraphs if p.text)
        has_vn, has_diac, diac_warn = check_vietnamese_diacritics(all_text)
        if has_vn and not has_diac:
            results.append(('Vietnamese diacritics', False, diac_warn))
        elif has_vn:
            results.append(('Vietnamese diacritics', True, 'OK'))

        # Check header/footer
        has_header = bool(s.header.paragraphs[0].text) if s.header else False
        has_footer = bool(s.footer.paragraphs[0].text) if s.footer else False
        results.append(('Header', has_header, s.header.paragraphs[0].text[:50] if has_header else 'empty'))
        results.append(('Footer', has_footer, 'present' if has_footer else 'empty'))

    except Exception as e:
        results.append(('File opens', False, str(e)))
    return results


def verify_xlsx(path):
    """Verify XLSX file."""
    from openpyxl import load_workbook
    results = []
    try:
        wb = load_workbook(path)
        results.append(('File opens', True, ''))

        # Sheets
        results.append(('Sheets', True, f'{len(wb.sheetnames)}: {", ".join(wb.sheetnames)}'))

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            row_count = ws.max_row
            col_count = ws.max_column
            results.append((f'[{sheet_name}] Size', True, f'{row_count} rows x {col_count} cols'))

            # Check formulas
            formula_count = 0
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith('='):
                        formula_count += 1
            # Formulas: pass if present, warn (not fail) if absent
            # Some sheets are text-only analysis — no formulas expected
            if formula_count > 0:
                results.append((f'[{sheet_name}] Formulas', True, f'{formula_count} formulas'))
            else:
                results.append((f'[{sheet_name}] Formulas', True, '0 formulas (text-only sheet?)'))

            # Check if header row has fill
            header_styled = False
            if ws.max_row > 0:
                first_cell = ws.cell(1, 1)
                if first_cell.fill and first_cell.fill.start_color and first_cell.fill.start_color.rgb:
                    header_styled = True
            results.append((f'[{sheet_name}] Styled', header_styled, ''))

            # Check freeze panes
            results.append((f'[{sheet_name}] Freeze panes', ws.freeze_panes is not None,
                          str(ws.freeze_panes or 'none')))

    except Exception as e:
        results.append(('File opens', False, str(e)))
    return results


def verify_pdf(path):
    """Verify PDF file."""
    results = []
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        page_count = len(reader.pages)
        results.append(('File opens', True, ''))
        results.append(('Pages', page_count > 0, f'{page_count} pages'))

        # Metadata (warnings, not errors — many valid PDFs lack metadata)
        meta = reader.metadata
        if meta:
            has_title = bool(meta.title)
            has_creator = bool(meta.creator)
            results.append(('Title', True, meta.title if has_title else '(none — consider adding)'))
            results.append(('Creator', True, meta.creator if has_creator else '(none — consider adding)'))

        # Extract text from first content page
        text = ''
        for page in reader.pages[:3]:
            text += page.extract_text() or ''

        results.append(('Has text', len(text) > 10, f'{len(text)} chars'))

        # Check Vietnamese characters
        vn_chars = set('àáảãạăắằẳẵặâấầẩẫậèéẻẽẹêếềểễệìíỉĩịòóỏõọôốồổỗộơớờởỡợùúủũụưứừửữựỳýỷỹỵđ')
        has_vn = any(c in text.lower() for c in vn_chars)
        if has_vn:
            # Check for replacement chars (font missing)
            has_tofu = '□' in text or '\ufffd' in text
            results.append(('Vietnamese text', True, 'found'))
            results.append(('Vietnamese renders', not has_tofu,
                          'OK' if not has_tofu else 'TOFU □ detected — font missing!'))
        else:
            results.append(('Vietnamese text', True, 'not present (English doc)'))

    except Exception as e:
        results.append(('File opens', False, str(e)))
    return results


def verify_pptx(path):
    """Verify PPTX file."""
    results = []
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slide_count = len(prs.slides)
        results.append(('File opens', True, ''))
        results.append(('Slides', slide_count > 0, f'{slide_count} slides'))

        # Check slide size (16:9 = 12192000 x 6858000 EMU or 13.333 x 7.5 inches)
        w = prs.slide_width
        h = prs.slide_height
        is_16_9 = abs(w / h - 16/9) < 0.1
        results.append(('Aspect ratio', True, '16:9' if is_16_9 else f'{w}x{h} EMU'))

        # Check fonts used
        fonts = set()
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
        results.append(('Fonts', True, ', '.join(fonts) or 'default'))

        # Check for non-web-safe fonts (warning)
        web_safe = {'Arial', 'Verdana', 'Tahoma', 'Georgia', 'Times New Roman',
                    'Trebuchet MS', 'Courier New', 'Helvetica', 'Impact'}
        non_safe = fonts - web_safe
        if non_safe:
            results.append(('Web-safe fonts', False,
                          f'Non-web-safe: {", ".join(non_safe)} \u2014 may not render on all machines'))

        # Check Vietnamese diacritics
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        all_text.append(para.text)
        full_text = ' '.join(all_text)
        has_vn, has_diac, diac_warn = check_vietnamese_diacritics(full_text)
        if has_vn and not has_diac:
            results.append(('Vietnamese diacritics', False, diac_warn))
        elif has_vn:
            results.append(('Vietnamese diacritics', True, 'OK'))

        # Check shapes per slide
        for i, slide in enumerate(prs.slides):
            shape_count = len(slide.shapes)
            results.append((f'Slide {i+1} shapes', True, f'{shape_count} shapes'))

    except Exception as e:
        results.append(('File opens', False, str(e)))
    return results


def verify_template_fill(original_path, output_path):
    """Compare original template vs filled output.
    Checks that structure is preserved after content replacement.
    Usage: python3 verify.py <output_path> --template <original_path>
    """
    ext = os.path.splitext(original_path)[1].lower()
    results = []

    if ext == '.docx':
        from docx import Document
        orig = Document(original_path)
        out = Document(output_path)
        # Compare paragraph count
        orig_paras = len(orig.paragraphs)
        out_paras = len(out.paragraphs)
        results.append(('Paragraph count', abs(orig_paras - out_paras) <= 2,
                        f'orig={orig_paras}, output={out_paras}'))
        # Compare table count
        results.append(('Table count', len(orig.tables) == len(out.tables),
                        f'orig={len(orig.tables)}, output={len(out.tables)}'))
        # Compare section count
        results.append(('Section count', len(orig.sections) == len(out.sections),
                        f'orig={len(orig.sections)}, output={len(out.sections)}'))
        # Compare heading count
        orig_h = len([p for p in orig.paragraphs if p.style.name.startswith('Heading')])
        out_h = len([p for p in out.paragraphs if p.style.name.startswith('Heading')])
        results.append(('Heading count', orig_h == out_h,
                        f'orig={orig_h}, output={out_h}'))

    elif ext == '.xlsx':
        from openpyxl import load_workbook
        orig = load_workbook(original_path)
        out = load_workbook(output_path)
        # Compare sheet names
        results.append(('Sheet names', orig.sheetnames == out.sheetnames,
                        f'orig={orig.sheetnames}, output={out.sheetnames}'))
        # Compare dimensions per sheet
        for name in orig.sheetnames:
            if name in out.sheetnames:
                o_ws = orig[name]
                n_ws = out[name]
                results.append((f'[{name}] Rows', o_ws.max_row == n_ws.max_row,
                                f'orig={o_ws.max_row}, output={n_ws.max_row}'))
                results.append((f'[{name}] Cols', o_ws.max_column == n_ws.max_column,
                                f'orig={o_ws.max_column}, output={n_ws.max_column}'))

    elif ext == '.pptx':
        from pptx import Presentation
        orig = Presentation(original_path)
        out = Presentation(output_path)
        results.append(('Slide count', len(orig.slides) == len(out.slides),
                        f'orig={len(orig.slides)}, output={len(out.slides)}'))

    else:
        results.append(('Template compare', False, f'Unsupported format: {ext}'))

    return results


def main():
    if len(sys.argv) < 2:
        print('Usage: python3 verify.py <file_path>')
        print('       python3 verify.py <output_path> --template <original_path>')
        sys.exit(1)

    path = sys.argv[1]
    template_path = None
    if '--template' in sys.argv:
        idx = sys.argv.index('--template')
        if idx + 1 < len(sys.argv):
            template_path = sys.argv[idx + 1]

    if not os.path.exists(path):
        print(f'❌ File not found: {path}')
        sys.exit(1)

    size = os.path.getsize(path)
    if size == 0:
        print(f'❌ File is empty: {path}')
        sys.exit(1)

    ext = os.path.splitext(path)[1].lower()
    print(f'📄 Verifying: {path} ({size:,} bytes)')
    print(f'   Format: {ext}')
    print()

    verifiers = {
        '.docx': verify_docx,
        '.xlsx': verify_xlsx,
        '.pdf': verify_pdf,
        '.pptx': verify_pptx,
    }

    if ext not in verifiers:
        print(f'❌ Unsupported format: {ext}')
        sys.exit(1)

    results = verifiers[ext](path)

    # Template comparison (if --template provided)
    if template_path:
        if os.path.exists(template_path):
            print(f'\n📝 Comparing with template: {template_path}')
            template_results = verify_template_fill(template_path, path)
            results.extend(template_results)
        else:
            print(f'\n⚠️  Template file not found: {template_path}')

    passed = 0
    failed = 0
    for name, ok, detail in results:
        icon = '✅' if ok else '❌'
        detail_str = f' — {detail}' if detail else ''
        print(f'  {icon} {name}{detail_str}')
        if ok:
            passed += 1
        else:
            failed += 1

    print()
    print(f'Result: {passed} passed, {failed} failed')
    if failed > 0:
        print('⚠️  Fix issues above before delivering file to user.')
        sys.exit(1)
    else:
        print('✅ All checks passed.')


if __name__ == '__main__':
    main()
